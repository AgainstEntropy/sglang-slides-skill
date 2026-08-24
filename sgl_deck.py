# SPDX-License-Identifier: Apache-2.0
"""sgl_deck — python-pptx helpers that encode the SGLang / RadixArk slide design system.

Import from a build script (PEP 723 + `uv run`):

    # /// script
    # requires-python = ">=3.10"
    # dependencies = ["python-pptx>=1.0", "qrcode[pil]"]
    # ///
    import sys, pathlib
    sys.path.insert(0, str(pathlib.Path(__file__).resolve().parents[N] / "sglang-slides-skill"))
    from sgl_deck import Deck, ACC, INK, ...

All coordinates are inches on a 13.333 x 7.5 (16:9) canvas.
Design tokens and layout rules: see design-system.md next to this file.
"""

from __future__ import annotations

import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---- design tokens (from RDXA happyhour deck) -------------------------------
INK = "1A1A1A"    # primary text
SUB = "555555"    # secondary text
MUT = "888888"    # muted / captions
LINE = "DDDDDD"   # hairlines, card borders
PANEL = "F5F5F5"  # card / panel fill
ACC = "C0492C"    # accent: eyebrows, key numbers, highlights
ACC2 = "2A7A6A"   # secondary accent: positive / quality / "safe"
TINT = "FBEFEA"   # light accent tint (lossy / warm)
TINT2 = "EAF3F0"  # light secondary tint (lossless / cool)
WHITE = "FFFFFF"

F = "Calibri"        # body + titles (safe-list font)
FM = "Courier New"   # code

W, H = 13.333, 7.5   # LAYOUT_WIDE canvas
MX = 0.6             # side margin
CW = W - 2 * MX      # content width (12.133)

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


class Deck:
    """A 16:9 deck with the SGLang design system baked in."""

    def __init__(self, author: str = "", title: str = ""):
        self.prs = Presentation()
        self.prs.slide_width = Inches(W)
        self.prs.slide_height = Inches(H)
        if author:
            self.prs.core_properties.author = author
        if title:
            self.prs.core_properties.title = title
        self._page = 0
        self._tmp: list[str] = []

    # -- slide ----------------------------------------------------------------
    def slide(self, page_num: bool = True):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = _rgb(WHITE)
        self._page += 1
        if page_num and self._page > 1:
            self.text(s, str(self._page), 12.7, 7.08, 0.45, 0.3, size=10, color=MUT, align="right")
        return s

    # -- text primitives --------------------------------------------------------
    def para(self, s, paragraphs, x, y, w, h, *, size=12.5, font=F, color=INK,
             align="left", valign="top", space_after=8, wrap=True):
        """paragraphs: list of paragraphs; each is a list of (text, style) runs.

        style keys: bold, italic, color, size, font, spc (char-spacing pt), bullet
        (bullet is a paragraph property — set it on the first run of the paragraph).
        """
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = _ANCHOR[valign]
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        for i, runs in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = _ALIGN[align]
            p.space_after = Pt(space_after)
            if runs and runs[0][1].get("bullet"):
                self._bullet(p)
            for text, st in runs:
                r = p.add_run()
                r.text = text
                r.font.name = st.get("font", font)
                r.font.size = Pt(st.get("size", size))
                r.font.bold = bool(st.get("bold", False))
                r.font.italic = bool(st.get("italic", False))
                r.font.color.rgb = _rgb(st.get("color", color))
                if st.get("spc"):
                    r.font._rPr.set("spc", str(int(st["spc"] * 100)))
        return tb

    @staticmethod
    def _bullet(p):
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "171450")   # 0.19" hanging indent
        pPr.set("indent", "-171450")
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
        pPr.append(buFont)
        pPr.append(buChar)

    def text(self, s, txt, x, y, w, h, *, size=12.5, bold=False, italic=False,
             color=INK, font=F, align="left", valign="top", spc=None, space_after=8, wrap=True):
        st = {"bold": bold, "italic": italic, "color": color, "size": size, "font": font}
        if spc:
            st["spc"] = spc
        paras = [[(line, dict(st))] for line in txt.split("\n")]
        return self.para(s, paras, x, y, w, h, size=size, font=font, color=color,
                         align=align, valign=valign, space_after=space_after, wrap=wrap)

    def rich(self, s, runs, x, y, w, h, *, size=12.5, align="left", valign="top", space_after=8):
        """One flow of styled runs; a run style {'br': True} starts a new paragraph after it."""
        paras, cur = [], []
        for text, st in runs:
            cur.append((text, st))
            if st.get("br"):
                paras.append(cur)
                cur = []
        if cur:
            paras.append(cur)
        return self.para(s, paras, x, y, w, h, size=size, align=align, valign=valign, space_after=space_after)

    # -- design-system components ----------------------------------------------
    def eyebrow(self, s, txt):
        self.text(s, txt, MX, 0.38, CW, 0.3, size=12, bold=True, color=ACC, spc=3)

    def title(self, s, txt, size=28):
        self.text(s, txt, MX, 0.66, CW, 0.75, size=size, bold=True, color=INK)

    def bullets(self, s, items, x, y, w, h, *, size=12.5, color=INK, space_after=8):
        paras = [[(t, {"bullet": True, "color": color, "size": size})] for t in items]
        self.para(s, paras, x, y, w, h, size=size, space_after=space_after)

    def lead_bullets(self, s, items, x, y, w, h, *, size=12.5, space_after=10):
        """items: [(bold lead, gray rest), ...] — the 'term — explanation' list pattern."""
        paras = [[(lead, {"bullet": True, "bold": True, "color": INK, "size": size}),
                  (rest, {"color": SUB, "size": size})] for lead, rest in items]
        self.para(s, paras, x, y, w, h, size=size, space_after=space_after)

    def stat_row(self, s, stats, *, y=1.4, x=MX, col_w=None, num_size=34, cap_size=11, color=ACC):
        """stats: [(number, caption), ...] — the big-number row."""
        cw = col_w or (CW / len(stats)) - 0.18
        step = col_w + 0.18 if col_w else CW / len(stats)
        for i, (num, cap) in enumerate(stats):
            self.text(s, num, x + i * step, y, cw, num_size / 50, size=num_size, bold=True, color=color)
            self.text(s, cap, x + i * step, y + num_size / 50 + 0.02, cw, 0.55, size=cap_size, color=SUB)

    # -- shapes ------------------------------------------------------------------
    def card(self, s, x, y, w, h, *, fill=PANEL, line=LINE, dash=False, radius=0.07):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.adjustments[0] = min(radius / min(w, h), 0.5)
        self._style(sh, fill, line, 1.0, dash)
        return sh

    def box(self, s, x, y, w, h, fill, txt=None, *, text_color=WHITE, size=10.5, line=None,
            dash=False, line_w=1.0, align="center", bold=False):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        self._style(sh, fill, line or fill, line_w, dash)
        if txt:
            self.text(s, txt, x, y, w, h, size=size, color=text_color, align=align,
                      valign="middle", bold=bold, space_after=0)
        return sh

    def ellipse(self, s, x, y, w, h, fill):
        sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        self._style(sh, fill, fill, 1.0, False)
        return sh

    def tri(self, s, x, y, w, h, fill, *, flip_h=False):
        sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        self._style(sh, fill, fill, 1.0, False)
        if flip_h:
            sh._element.spPr.xfrm.set("flipH", "1")
        return sh

    def line(self, s, x1, y1, x2, y2, *, color=MUT, width=1.5, dash=False):
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        self._no_theme_fx(ln)
        ln.line.color.rgb = _rgb(color)
        ln.line.width = Pt(width)
        if dash:
            self._dash(ln.line)
        return ln

    def arrow(self, s, x1, y1, x2, y2, *, color=MUT, width=1.5, dash=False):
        ln = self.line(s, x1, y1, x2, y2, color=color, width=width, dash=dash)
        lnEl = ln.line._get_or_add_ln()
        lnEl.append(lnEl.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
        return ln

    @staticmethod
    def _dash(line_format):
        lnEl = line_format._get_or_add_ln()
        lnEl.append(lnEl.makeelement(qn("a:prstDash"), {"val": "dash"}))

    @staticmethod
    def _no_theme_fx(sh):
        """Zero the theme effect reference — LibreOffice renders it as a drop
        shadow even when an empty <a:effectLst/> is present."""
        st = sh._element.find(qn("p:style"))
        if st is not None:
            er = st.find(qn("a:effectRef"))
            if er is not None:
                er.set("idx", "0")

    @staticmethod
    def _style(sh, fill, line, line_w, dash):
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(fill)
        sh.line.color.rgb = _rgb(line)
        sh.line.width = Pt(line_w)
        if dash:
            Deck._dash(sh.line)
        sh.shadow.inherit = False
        Deck._no_theme_fx(sh)
        tf = sh.text_frame
        tf.word_wrap = True

    # -- table ---------------------------------------------------------------------
    def table(self, s, rows, x, y, col_w, row_h, *, size=12, border=LINE, border_pt=0.5, margin=0.06):
        """rows: list of rows; a cell is text or (text, opts).

        opts: bold, color, fill, align. First row is NOT special — style it yourself.
        """
        n_r, n_c = len(rows), len(rows[0])
        shp = s.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(sum(col_w)), Inches(row_h * n_r))
        tbl = shp.table
        # kill the default banded style
        tblPr = tbl._tbl.tblPr
        for attr in ("firstRow", "bandRow", "firstCol", "lastRow", "lastCol", "bandCol"):
            tblPr.set(attr, "0")
        for i, wid in enumerate(col_w):
            tbl.columns[i].width = Inches(wid)
        for i in range(n_r):
            tbl.rows[i].height = Inches(row_h)
        for ri, row in enumerate(rows):
            for ci, cell_def in enumerate(row):
                txt, opts = cell_def if isinstance(cell_def, tuple) else (cell_def, {})
                cell = tbl.cell(ri, ci)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                for m in ("margin_left", "margin_right"):
                    setattr(cell, m, Inches(margin))
                for m in ("margin_top", "margin_bottom"):
                    setattr(cell, m, Inches(0.02))
                self._cell_borders(cell, border, border_pt)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(opts.get("fill", WHITE))
                p = cell.text_frame.paragraphs[0]
                p.alignment = _ALIGN[opts.get("align", "left")]
                r = p.add_run()
                r.text = txt
                r.font.name = F
                r.font.size = Pt(opts.get("size", size))
                r.font.bold = bool(opts.get("bold", False))
                r.font.color.rgb = _rgb(opts.get("color", INK))
        return tbl

    @staticmethod
    def _cell_borders(cell, color, pt):
        tcPr = cell._tc.get_or_add_tcPr()
        for i, tag in enumerate(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
            ln = tcPr.find(qn(tag))
            if ln is None:
                ln = tcPr.makeelement(qn(tag), {})
                tcPr.insert(i, ln)  # ln* must precede the fill element
            ln.set("w", str(int(Pt(pt))))  # a:ln width is in EMU
            for child in list(ln):
                ln.remove(child)
            fillEl = ln.makeelement(qn("a:solidFill"), {})
            clr = ln.makeelement(qn("a:srgbClr"), {"val": color})
            fillEl.append(clr)
            ln.append(fillEl)

    # -- extras ------------------------------------------------------------------
    def code_panel(self, s, code, x, y, w, h, *, size=9.5, fill=INK, color="EEEEEE", pad=0.2):
        self.card(s, x, y, w, h, fill=fill, line=fill, radius=0.07)
        self.text(s, code, x + pad, y + pad * 0.75, w - 2 * pad, h - 1.5 * pad,
                  size=size, color=color, font=FM, space_after=0)

    def image(self, s, path, x, y, w=None, h=None):
        kw = {}
        if w:
            kw["width"] = Inches(w)
        if h:
            kw["height"] = Inches(h)
        return s.shapes.add_picture(str(path), Inches(x), Inches(y), **kw)

    def qr(self, s, url, x, y, w, *, caption=None, caption_w=None):
        """Generate an ink-colored QR for `url` and place it."""
        import qrcode

        q = qrcode.QRCode(border=1, box_size=10, error_correction=qrcode.constants.ERROR_CORRECT_M)
        q.add_data(url)
        q.make(fit=True)
        img = q.make_image(fill_color="#1A1A1A", back_color="white")
        f = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        img.save(f.name)
        self._tmp.append(f.name)
        pic = self.image(s, f.name, x, y, w=w, h=w)
        if caption:
            cw = caption_w or w + 1.4
            self.text(s, caption, x + w / 2 - cw / 2, y + w + 0.08, cw, 0.5,
                      size=10.5, color=MUT, align="center")
        return pic

    def save(self, path):
        self.prs.save(str(path))
        for f in self._tmp:
            Path(f).unlink(missing_ok=True)
        print("wrote", path)
